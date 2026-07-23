# -*- coding: utf-8 -*-
"""예인 시티팝 6곡 가사 평가 보고서 PPTX 생성 (곡당 1페이지, 4관점 평가)."""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

BG = RGBColor(0x12, 0x12, 0x1A); CARD = RGBColor(0x1E, 0x1E, 0x2A)
MAGENTA = RGBColor(0xE9, 0x5B, 0xB5); CYAN = RGBColor(0x5B, 0xC8, 0xE9)
WHITE = RGBColor(0xEE, 0xEE, 0xF2); GREY = RGBColor(0x9A, 0x9A, 0xAC)
AMBER = RGBColor(0xE9, 0xB4, 0x5B); GREEN = RGBColor(0x8B, 0xE9, 0x5B)
FONT = "맑은 고딕"

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def bgfill(s):
    r = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG; r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2, r._element)

def tb(s, x, y, w, h, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT, ls=1.0):
    box = s.shapes.add_textbox(x, y, w, h); tf = box.text_frame; tf.word_wrap = True
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln; p.alignment = align; p.line_spacing = ls
        for r in p.runs:
            r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    return box

def card(s, x, y, w, h, accent):
    c = s.shapes.add_shape(5, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = CARD
    c.line.color.rgb = accent; c.line.width = Pt(1.2); c.shadow.inherit = False

def stars(n):
    full = int(n); half = (n - full) >= 0.5
    return "★" * full + ("☆" if half else "") + f"  {n:.1f}"

SONGS = [
 dict(no="A01", title="네온 (Neon)", concept="이별 후 일정하지 않은 마음 — 네온은 켜져 있는 게 아니라 깜빡인다",
  info="세트 A(도시의 밤) · 마이너 멜랑콜리 시티팝 · 98 BPM · 후렴 sax\n상태: Suno 확정본 완료 · 리릭비디오 발행 · LRC 싱크 완료",
  lit="'켜진/꺼진' 한 단어만 뒤집는 극성 대구(V2-Bridge)가 곡의 척추. 마음의 진폭을 설명하지 않고 구조로 실연한다. 이름 모티프가 지움(V1)-보임(Bridge)-쫓음(Outro)으로 3박자 회수되는 결말은 시적 완결성이 높다. 감정을 신호등과 간판으로 번역하는 절제가 세트 전체의 기준점.",
  crit="훅 '깜빡거리던 건가'(7음절)의 절벽 낙차는 강하지만 세 후렴 동일 종결의 단조로움 리스크는 남는다. '너' 심기를 프리코러스 2행으로 최소화한 판단은 세련되나 이별 서사 가독성은 여전히 경계선 — 초견에 '도시 소진'으로 오독될 여지가 있다.",
  comp="98BPM에 13/12/12/12/7 후렴 낙차 — 마지막 7음절에 멜로디 훅이 앉을 자리가 명확하고 늘여 부를 공간도 있다. 프리코러스 '쪽으로/쪽으로' 병렬은 후렴 진입 가속 장치. 편곡 지시(sax 훅, 게이트 리버브)가 구체적이라 프로덕션 방향이 흔들리지 않는다.",
  pub="첫 청취에 '네온, 깜빡' 이미지가 바로 박힌다. 훅이 짧아 흥얼거리기 좋고, 2절쯤 이별 노래임을 알아채는 재미가 있다. 밤 드라이브 플레이리스트 1순위 감.",
  s=(5.0, 4.0, 4.5, 4.0)),
 dict(no="A02", title="워커홀릭", concept="탈출 욕구마저 사라진 소진 — 버티게 하는 건 캔맥주 하나",
  info="세트 A · cold minimal 시티팝(no sax, 신스솔로) · 100 BPM\n상태: Suno 완료 — 시티팝예인 2 + 예인확정본 2 (A/B 대기) · 키보드 스틸 승인",
  lit="감정어 없이 '어제 한 일이 오늘 다시 와 있고' 같은 반복 사실만으로 소진을 실연한다. '알람까지 네 시간 반'의 구체 숫자, '한 캔만 더'의 조사 '만' — 미시 언어 운용이 백미. 마지막 후렴의 장소 이동(사무실-집)은 형식이 주제(도망칠 곳 없음)를 수행하는 드문 사례.",
  crit="볼빨간사춘기 동명곡과의 차별화(일탈 아닌 소진)는 명확히 성립. 다만 브릿지 캔맥주가 유일한 온기라 정서가 단색으로 흐를 위험. 아웃트로가 '설명'으로 미끄러질 경계를 '~였으면 해서'의 닫히지 않는 어미로 아슬하게 지켜냈다.",
  comp="프리코러스를 제거한 단조 골격 + FM 신스 솔로 삽입 — 사운드 정체성이 뚜렷하다. 기계적 반복 그루브에 deadpan 창법 지시는 주제와 정합. 관건은 Suno가 단조 구조에서 훅을 안착시키는가.",
  pub="'끝내면 집에 가려고 했는데' — 직장인 공감 직격. 밈과 짤로 소비될 훅. 다만 위로 없이 차갑게 끝나 반복청취는 취향을 탄다.",
  s=(4.5, 4.0, 4.0, 4.0)),
 dict(no="A03", title="회식", concept="웃는 타이밍을 외워서 웃는다 — 여럿 속의 소외",
  info="세트 A · groovy 멜랑콜리 시티팝(북적-희박) · 105 BPM\n상태: Suno 완료(2버전) · 씬 스틸 리워크 발행(배시시+회전블러)",
  lit="'웃는 타이밍은 이제 알겠어' 한 줄이 사회적 가면을 요약한다. 브릿지의 거울 앞 '어릴 적 아빠의 표정'은 소외의 세대 전승을 지나치듯 던지는 문학적 정점. 종결 '문을 열지. 직장인. 간다.'의 단문 3연타는 산문시적 처리로 체념을 리듬화한다.",
  crit="취기 장치('내가 흔들리는지 사람들이 흔들리는 건지')로 시점을 흔든 설계가 좋다. 다만 '숙취해소제/꿀꺽'의 즉물성은 신선함과 개그 사이 아슬한 줄타기 — 보컬 톤이 무게를 못 잡으면 우스개로 미끄러진다.",
  comp="파티 그루브가 후반에 희박해지는 다이내믹 서사 지시가 가사 구조와 정확히 맞물린다. 후렴 픽업 '그 와중에도'는 그루브에 태우기 좋은 진입구. '꿀꺽' 한 단어 행은 리듬 브레이크 포인트로 재미 요소.",
  pub="회식 다녀온 날 귀가길 재생 각. '맞장구는 몸이 먼저 쳐'는 짤 예감. 아빠 표정 라인에서 훅 먹먹해지는 감정 낙차가 좋다.",
  s=(4.5, 3.5, 4.0, 4.5)),
 dict(no="A04", title="번화가", concept="저 불빛은 전부 누군가의 것 — 화려함 속 존재감 없음",
  info="세트 A · bright glossy 업템포 시티팝(슬랩베이스, 브라스) · 110 BPM\n상태: Suno 완료(2버전) · 씬 스틸 발행(인파 모션블러)",
  lit="'간판은 저마다 이름이 있고'로 열어 '저기 있었던 적이 있었나 싶어'로 닫는, 익명 존재의 증명 실패담. 전단지 받기, 부딪힘, 프레임 밖으로 비켜주기 등 도시의 미세 상호작용만으로 소외를 그린다. '비켜 주는 건 어렵지 않지'의 자조가 곡에서 가장 아픈 축.",
  crit="코러스 선행 구조는 A세트 내 골격 다양화에 기여. 다만 훅의 질문형 종결('뭘 보고 있나')이 네온, 신호대기의 질문 어미와 겹쳐 세트를 연속 청취할 때 변별력이 약해질 우려가 있다.",
  comp="110BPM 화려 편성(네온 브라스, 슬랩)과 공허한 가사의 대비 — 시티팝 정석 설계. 후렴으로 시작하니 도입 3초 훅이 가능하고 스트리밍 스킵 방어에 유리하다.",
  pub="사운드에 몸 흔들다 가사를 곱씹으면 쓸쓸해지는 이중 재미. 지하철 입구에서 뒤돌아보는 마지막 장면의 공감도가 높다.",
  s=(4.0, 3.5, 4.5, 4.0)),
 dict(no="A05", title="신호대기", concept="차는 한 대도 안 오는데 — 아무도 안 건너는 빨간불",
  info="세트 A · tense minimal 시티팝(stop-and-go, 중간 정적) · 100 BPM\n상태: Suno 완료(2버전) · 씬 스틸 생성 중 · A05-A06-A07 3부작의 1편",
  lit="텅 빈 도로의 빨간불 앞 군중 — 내면화된 규율의 알레고리를 일상 스냅 하나로 세운다. '어길 수가 없어. 보는 눈이 많거든'의 마침표와 구어 종결이 억압의 자발성을 폭로한다. 아웃트로 '기다린 건 같이지만/향하는 곳이 다른데'는 군중이라는 공동체의 역설로 닫는 좋은 결구.",
  crit="[Break] 정적 구간 — 형식이 '멈춤'이라는 주제를 수행하는 설계는 높이 살 만하다. 다만 스무 명, 열둘 등 관찰-숫자 나열이 리스트에 머물 위험이 있어 후렴 반복이 정서를 끌어올릴지가 관건.",
  comp="코일드 stop-and-go 그루브와 중간 무음 브레이크는 편곡 난이도이자 차별화 기회. 후렴이 4행으로 짧아 반복 훅 설계가 필요하다. 스타카토 지시가 명확해 리듬 프로덕션 방향은 선명.",
  pub="'차는 한 대도 안 오는데' — 전 국민이 겪은 순간의 소환력. 후렴이 짧고 중독성 있다. 다만 수수한 곡이라 첫 귀 임팩트는 세트 중 중간.",
  s=(4.0, 4.0, 3.5, 3.5)),
 dict(no="B01", title="주말에", concept="사별처럼 들리지만 전부 사실 — 자취생의 된장찌개",
  info="세트 B(가족, 집) · 따뜻한 어쿠스틱 시티팝(플루겔혼) · 90 BPM\n상태: Suno 완료(2버전) · 실화 기반(첫 술 마신 날 엄마에게 전화)",
  lit="모든 행이 사별로 읽히도록 설계됐지만 거짓 문장이 단 하나도 없다 — 중의성 하나로 곡 전체를 끄는 구성적 재치. 금지어 원칙('보고 싶다' 부재)이 미끼를 살리고, '그래 밥 해둘게, 하고/엄마가 먼저 끊네'의 무심한 종결이 반전과 정서를 동시에 회수한다. 웃기고 나서 더 찡한, 세트에서 가장 영리한 결말.",
  crit="'엄마 없이'의 반복은 오독 설계의 핵이자 윤리적 아슬함 — 사별 마케팅으로 비칠 위험이 있으나 허위 문장이 없다는 방어선이 유효하다. 반전 후 전곡이 다르게 읽히는 재청취 가치는 세트 최고 수준.",
  comp="장례처럼 무겁게 쌓다 마지막에 '툭' 내려놓는 다이내믹 — 아웃트로에서 반주를 걷어내는 처리가 킬 포인트가 될 것. sax 대신 플루겔혼 지정 등 편성 온도 감각이 정확하다.",
  pub="1절에서 울고 아웃트로에서 피식 — 감정 롤러코스터의 쾌감. '청양고추 팍팍' 디테일이 다 살린다. 자취생 필청, SNS 사연 태그 예감. 대중 소구력은 6곡 중 최고.",
  s=(4.5, 4.0, 4.0, 5.0)),
]

# 표지
s = prs.slides.add_slide(blank); bgfill(s)
tb(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.0), "예인(Ye_in) 시티팝 프로젝트", 40, WHITE, True)
tb(s, Inches(0.9), Inches(3.0), Inches(11.5), Inches(0.8), "가사 평가 보고서 — 6곡", 28, MAGENTA, True)
bar = s.shapes.add_shape(1, Inches(0.9), Inches(3.95), Inches(3.2), Pt(3))
bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background(); bar.shadow.inherit = False
tb(s, Inches(0.9), Inches(4.25), Inches(11.5), Inches(1.6),
   "평가 관점: 문학평론가 · 비평가 · 대중음악작곡가 · 대중\n대상: A01 네온 · A02 워커홀릭 · A03 회식 · A04 번화가 · A05 신호대기 · B01 주말에\n2026-07-23 · 작성: Claude (예인 프로젝트 세션)", 15, GREY, ls=1.35)

LABELS = [("문학평론가", MAGENTA), ("비평가", CYAN), ("대중음악작곡가", AMBER), ("대중", GREEN)]
for sg in SONGS:
    s = prs.slides.add_slide(blank); bgfill(s)
    tb(s, Inches(0.55), Inches(0.28), Inches(1.3), Inches(0.6), sg["no"], 26, CYAN, True)
    tb(s, Inches(1.55), Inches(0.22), Inches(8.6), Inches(0.7), sg["title"], 30, WHITE, True)
    tb(s, Inches(1.58), Inches(0.82), Inches(11.2), Inches(0.45), sg["concept"], 13.5, MAGENTA)
    tb(s, Inches(0.58), Inches(1.32), Inches(12.2), Inches(0.75), sg["info"], 11.5, GREY, ls=1.25)
    keys = ["lit", "crit", "comp", "pub"]
    W = Inches(6.02); H = Inches(2.42); X0 = Inches(0.55); Y0 = Inches(2.18)
    GX = Inches(0.18); GY = Inches(0.16)
    for i, (lab, accent) in enumerate(LABELS):
        x = X0 + (W + GX) * (i % 2); y = Y0 + (H + GY) * (i // 2)
        card(s, x, y, W, H, accent)
        tb(s, x + Inches(0.22), y + Inches(0.10), Inches(3.6), Inches(0.4), lab, 14, accent, True)
        tb(s, x + W - Inches(2.35), y + Inches(0.12), Inches(2.15), Inches(0.4), stars(sg["s"][i]), 12.5, accent, True, PP_ALIGN.RIGHT)
        tb(s, x + Inches(0.22), y + Inches(0.52), W - Inches(0.44), H - Inches(0.66), sg[keys[i]], 11, WHITE, ls=1.18)

# 총평
s = prs.slides.add_slide(blank); bgfill(s)
tb(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.8), "총평", 30, WHITE, True)
summary = (
"세트 전체가 '감정을 명명하지 않고 사물과 행동으로 번역한다'는 하나의 시학으로 묶여 있으면서도, 곡마다 다른 구조 장치(극성 대구, 장소 이동 후렴, 단문 종결, 코러스 선행, 정적 브레이크, 중의성 반전)를 세워 골격 중복을 피했다 — 연작으로서의 설계 완성도가 높다.\n\n"
"강점: 1) 훅의 낙차 설계(네온의 7음절 절벽, 주말에의 '툭')가 일관되게 작동한다. 2) 미시 언어(조사 '만', 마침표 종결, 금지어 원칙) 운용이 프로 작사 수준. 3) 사운드 지시가 가사 주제와 정합해 Suno 프로덕션 방향이 흔들리지 않는다.\n\n"
"과제: 1) 질문형 어미('~건가/~있나')가 세트 전반에 반복되어 연속 청취 시 변별이 떨어진다 — A06·A07 작사 시 다른 종결 전략 권장. 2) A02·A05의 단색 정서는 편곡(신스 솔로, 브레이크) 완성도가 성패를 가른다. 3) B01의 오독 설계는 공개 시 '실화, 아무도 죽지 않음' 한 줄로 방어선을 명시할 것.\n\n"
"종합: 문학성은 A01·B01, 대중성은 B01·A03이 앞서고, 세트의 얼굴은 여전히 A01 네온이다. 다음 작업은 A06 막차·A07 첫차 가사 완성으로 3부작(A05-06-07)의 감정 곡선을 닫는 것."
)
tb(s, Inches(0.9), Inches(1.4), Inches(11.6), Inches(5.6), summary, 13.5, WHITE, ls=1.3)

os.makedirs("lyric_output", exist_ok=True)
out = "lyric_output/예인_가사평가보고서_6곡.pptx"
prs.save(out)
print("SAVED", out, os.path.getsize(out), "bytes, slides:", len(prs.slides._sldIdLst))
